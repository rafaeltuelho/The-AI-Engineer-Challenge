"""
Lightweight document processing and RAG functionality supporting PDF, Word (.docx), and PowerPoint (.pptx).
This version avoids heavy ML dependencies for Vercel deployment.
"""
import os
import tempfile
import uuid
from typing import List, Dict, Any, Optional
from pathlib import Path
import asyncio
import re
import mimetypes

import PyPDF2
import pdfplumber
import tiktoken
from docx import Document
from pptx import Presentation

import sys
import os
sys.path.append(os.path.join(os.path.dirname(__file__), '..'))

from aimakerspace.vectordatabase import VectorDatabase
from aimakerspace.openai_utils.embedding import EmbeddingModel
from aimakerspace.openai_utils.chatmodel import ChatOpenAI


class DocumentProcessor:
    """Lightweight document processor supporting PDF, Word (.docx), and PowerPoint (.pptx)."""
    
    def __init__(self):
        # Initialize tokenizer for chunking
        self.tokenizer = tiktoken.encoding_for_model("text-embedding-3-small")
        
    def process_document(self, file_path: str) -> Dict[str, Any]:
        """
        Process a document file and extract content with chunks.
        Supports PDF, Word (.docx), and PowerPoint (.pptx) files.
        
        Args:
            file_path: Path to the document file
            
        Returns:
            Dictionary containing processed content and metadata
        """
        try:
            # Check if file exists and is readable
            if not os.path.exists(file_path):
                raise Exception(f"Document file not found: {file_path}")
            
            # Check file size (limit to 50MB to prevent memory issues)
            file_size = os.path.getsize(file_path)
            if file_size > 50 * 1024 * 1024:  # 50MB limit
                raise Exception(f"Document file too large: {file_size / (1024*1024):.1f}MB (max 50MB)")
            
            # Detect file type
            file_type = self._detect_file_type(file_path)
            file_name = os.path.basename(file_path)
            
            print(f"Processing {file_type.upper()} document: {file_name} ({file_size / (1024*1024):.1f}MB)")
            
            # Extract text based on file type
            text_content = self._extract_text_by_type(file_path, file_type)
            
            if not text_content.strip():
                raise Exception(f"No text could be extracted from the {file_type} document")
            
            print(f"Extracted text length: {len(text_content)} characters")
            
            # Create chunks
            chunks = self._create_chunks(text_content)
            
            print(f"Created {len(chunks)} chunks")
            
            return {
                "full_text": text_content,
                "chunks": chunks,
                "chunk_count": len(chunks),
                "metadata": {
                    "file_path": file_path,
                    "file_name": file_name,
                    "file_type": file_type,
                    "file_size": file_size,
                    "processing_method": f"lightweight_{file_type}"
                }
            }
            
        except Exception as e:
            print(f"Error processing document: {str(e)}")
            raise Exception(f"Error processing document: {str(e)}")
    
    def _detect_file_type(self, file_path: str) -> str:
        """Detect the file type based on extension and MIME type."""
        file_extension = Path(file_path).suffix.lower()
        
        if file_extension in ['.pdf']:
            return 'pdf'
        elif file_extension in ['.docx', '.doc']:
            return 'docx'
        elif file_extension in ['.pptx', '.ppt']:
            return 'pptx'
        else:
            # Try MIME type detection as fallback
            mime_type, _ = mimetypes.guess_type(file_path)
            if mime_type:
                if 'pdf' in mime_type:
                    return 'pdf'
                elif 'wordprocessingml' in mime_type or 'msword' in mime_type:
                    return 'docx'
                elif 'presentationml' in mime_type or 'mspowerpoint' in mime_type:
                    return 'pptx'
            
            raise Exception(f"Unsupported file type: {file_extension}. Supported types: .pdf, .docx, .pptx")
    
    def _extract_text_by_type(self, file_path: str, file_type: str) -> str:
        """Extract text based on file type."""
        if file_type == 'pdf':
            return self._extract_text_pdf(file_path)
        elif file_type == 'docx':
            return self._extract_text_docx(file_path)
        elif file_type == 'pptx':
            return self._extract_text_pptx(file_path)
        else:
            raise Exception(f"Unsupported file type for extraction: {file_type}")
    
    def _extract_text_pdf(self, pdf_path: str) -> str:
        """Extract text from PDF using both PyPDF2 and pdfplumber for better results."""
        text_pypdf2 = self._extract_text_pypdf2(pdf_path)
        text_pdfplumber = self._extract_text_pdfplumber(pdf_path)
        
        # Use pdfplumber if it extracted more text, otherwise use PyPDF2
        if len(text_pdfplumber.strip()) > len(text_pypdf2.strip()):
            return text_pdfplumber
        else:
            return text_pypdf2
    
    def _extract_text_docx(self, docx_path: str) -> str:
        """Extract text from Word document (.docx)."""
        try:
            doc = Document(docx_path)
            text = ""
            
            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text += paragraph.text + "\n"
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text += " | ".join(row_text) + "\n"
            
            return text
        except Exception as e:
            print(f"Word document extraction failed: {e}")
            return ""
    
    def _extract_text_pptx(self, pptx_path: str) -> str:
        """Extract text from PowerPoint presentation (.pptx)."""
        try:
            prs = Presentation(pptx_path)
            text = ""
            
            for slide_num, slide in enumerate(prs.slides):
                # Add slide header
                text += f"Slide {slide_num + 1}:\n"
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
                    
                    # Extract text from tables
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                text += " | ".join(row_text) + "\n"
                
                text += "\n"  # Add spacing between slides
            
            return text
        except Exception as e:
            print(f"PowerPoint extraction failed: {e}")
            return ""
    
    def _extract_text_pypdf2(self, pdf_path: str) -> str:
        """Extract text using PyPDF2."""
        try:
            with open(pdf_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                # Safety check for large PDFs
                if len(pdf_reader.pages) > 100:
                    print(f"Warning: Large PDF detected ({len(pdf_reader.pages)} pages), processing first 50 pages only")
                    max_pages = 50
                else:
                    max_pages = len(pdf_reader.pages)
                
                text = ""
                for i, page in enumerate(pdf_reader.pages[:max_pages]):
                    try:
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + "\n"
                    except Exception as page_error:
                        print(f"Error extracting page {i}: {page_error}")
                        continue
                        
                return text
        except Exception as e:
            print(f"PyPDF2 extraction failed: {e}")
            return ""
    
    def _extract_text_pdfplumber(self, pdf_path: str) -> str:
        """Extract text using pdfplumber."""
        try:
            with pdfplumber.open(pdf_path) as pdf:
                # Safety check for large PDFs
                if len(pdf.pages) > 100:
                    print(f"Warning: Large PDF detected ({len(pdf.pages)} pages), processing first 50 pages only")
                    max_pages = 50
                else:
                    max_pages = len(pdf.pages)
                
                text = ""
                for i, page in enumerate(pdf.pages[:max_pages]):
                    try:
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + "\n"
                    except Exception as page_error:
                        print(f"Error extracting page {i}: {page_error}")
                        continue
                        
                return text
        except Exception as e:
            print(f"pdfplumber extraction failed: {e}")
            return ""
    
    def _create_chunks(self, text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
        """
        Create text chunks using simple token-based splitting.
        
        Args:
            text: Text to chunk
            chunk_size: Maximum tokens per chunk
            overlap: Number of tokens to overlap between chunks
        """
        if not text.strip():
            return []
        
        try:
            # Tokenize the text
            tokens = self.tokenizer.encode(text)
            
            # Safety check for very large texts
            if len(tokens) > 100000:  # Limit to prevent memory issues
                print(f"Warning: Large text detected ({len(tokens)} tokens), truncating...")
                tokens = tokens[:100000]
            
            chunks = []
            start = 0
            max_iterations = len(tokens) // (chunk_size - overlap) + 10  # Safety limit
            iteration_count = 0
            
            while start < len(tokens) and iteration_count < max_iterations:
                iteration_count += 1
                
                # Define the end of the current chunk
                end = min(start + chunk_size, len(tokens))
                
                # Extract chunk tokens
                chunk_tokens = tokens[start:end]
                
                # Decode back to text
                chunk_text = self.tokenizer.decode(chunk_tokens)
                
                # Clean up the chunk text
                chunk_text = self._clean_chunk_text(chunk_text)
                
                if chunk_text.strip():
                    chunks.append(chunk_text)
                
                # Move to the next chunk with overlap
                start = end - overlap
                
                # Safety check to prevent infinite loop
                if start >= end or start < 0:
                    break
                    
                # Additional safety check
                if start >= len(tokens):
                    break
            
            return chunks
            
        except Exception as e:
            print(f"Error in chunking: {e}")
            # Fallback to simple text splitting
            return self._create_simple_chunks(text, chunk_size=1000)
    
    def _create_simple_chunks(self, text: str, chunk_size: int = 1000) -> List[str]:
        """Fallback method for creating chunks using simple text splitting."""
        if not text.strip():
            return []
        
        # Split text into words
        words = text.split()
        chunks = []
        
        for i in range(0, len(words), chunk_size):
            chunk_words = words[i:i + chunk_size]
            chunk_text = ' '.join(chunk_words)
            if chunk_text.strip():
                chunks.append(chunk_text.strip())
        
        return chunks
    
    def _clean_chunk_text(self, text: str) -> str:
        """Clean and normalize chunk text."""
        # Remove excessive whitespace
        text = re.sub(r'\s+', ' ', text)
        # Remove leading/trailing whitespace
        text = text.strip()
        return text

    # Backward compatibility method
    def process_pdf(self, pdf_path: str) -> Dict[str, Any]:
        """
        Backward compatibility method for PDF processing.
        Use process_document() for new implementations.
        
        Args:
            pdf_path: Path to the PDF file
            
        Returns:
            Dictionary containing processed content and metadata
        """
        return self.process_document(pdf_path)


class RAGSystem:
    """RAG (Retrieval-Augmented Generation) system using aimakerspace library."""
    
    def __init__(self, api_key: str, model_name: str = "gpt-4o-mini"):
        self.api_key = api_key
        self.model_name = model_name
        self.embedding_model = EmbeddingModel(api_key=api_key)
        self.vector_db = VectorDatabase(embedding_model=self.embedding_model, api_key=api_key)
        self.chat_model = ChatOpenAI(model_name=model_name, api_key=api_key)
        self.documents = {}  # Store document metadata
        self.topic_explorer_system_message = """
        You are an educational study companion for middle school students learning Math, Science, or US History. 
        Your role is to help students understand topics from their class materials in a clear, friendly, and encouraging way. 
        Always explain ideas at the level of an elementary or middle school student, avoiding overly complex words. 

        When answering a question, always follow this structure:
        1. **Explanation:** a simple, clear explanation based on the provided context, using age-appropriate language.
        2. **Real-Life Example:** show how the idea connects to something in the student's everyday life.
        3. **Practice Activity:** create a short, fun challenge (problem to solve, small writing task, or drawing prompt) that helps the student practice.

        Do not give long essays. Be concise, supportive, and engaging, like a tutor who makes learning fun.
        """
        self.rag_system_message = """
        You are a helpful assistant that answers questions based on provided context.
        """

    async def index_document(self, document_id: str, chunks: List[str], metadata: Dict[str, Any]) -> None:
        """
        Index document chunks in the vector database.
        
        Args:
            document_id: Unique identifier for the document
            chunks: List of text chunks to index
            metadata: Document metadata
        """
        try:
            # Get embeddings for all chunks
            embeddings = await self.embedding_model.async_get_embeddings(chunks)
            
            # Store each chunk in the vector database
            for i, (chunk, embedding) in enumerate(zip(chunks, embeddings)):
                chunk_id = f"{document_id}_chunk_{i}"
                
                # Create metadata for this chunk
                chunk_metadata = {
                    "document_id": document_id,
                    "chunk_index": i,
                    "chunk_text": chunk,
                    **metadata
                }
                
                # Insert into vector database
                self.vector_db.insert(chunk_id, embedding, chunk_metadata)
            
            # Store document metadata
            self.documents[document_id] = {
                "chunk_count": len(chunks),
                "metadata": metadata
            }
            
        except Exception as e:
            raise Exception(f"Error indexing document: {str(e)}")
    
    def search_relevant_chunks(self, query: str, k: int = 3) -> List[Dict[str, Any]]:
        """
        Search for relevant chunks based on query.
        
        Args:
            query: Search query
            k: Number of top chunks to return
        
        Returns:
            List of relevant chunks with metadata
        """
        try:
            # Use search_with_metadata method that returns dictionaries directly
            results = self.vector_db.search_with_metadata(query, k=k)
            
            return results
            
        except Exception as e:
            raise Exception(f"Error searching chunks: {str(e)}")
    
    async def query_documents(self, query: str, k: int = 3, mode: str = "rag") -> str:
        """
        Query documents using RAG approach.
        
        Args:
            query: User query
            k: Number of chunks to retrieve
        
        Returns:
            Generated response based on retrieved chunks
        """
        try:
            # Search for relevant chunks
            relevant_chunks = self.search_relevant_chunks(query, k=k)
            
            if not relevant_chunks:
                return "I couldn't find any relevant information in the uploaded documents to answer your question."
            
            # Prepare context from chunks
            context_parts = []
            for chunk_data in relevant_chunks:
                chunk_text = chunk_data.get('chunk_text', '')
                if chunk_text:
                    context_parts.append(chunk_text)
            
            context = "\n\n".join(context_parts)
            
            # Create prompt based on mode
            if mode == "topic-explorer":
                # Topic Explorer mode with structured 3-part response
                rag_prompt = f"""You are a helpful and engaging study assistant for middle school students.
You will receive:
1. A student question
2. Context extracted from their study material (PDF, Word, or PowerPoint)

Your task:
- Use the context when possible, but also explain in a clear, simple way.
- Always return the answer in **three structured sections**:
  1. **Explanation** → Simple, student-friendly explanation of the topic.
  2. **Real-Life Example** → Show how the concept applies in daily life or something relatable.
  3. **Practice Activity** → Give a short challenge (question, drawing, or exercise) the student can do to practice.

Format your answer clearly with section headers.

---
Student Question:
{query}

Context:
{context}

Answer:"""
                system_message = self.topic_explorer_system_message
            else:
                # Regular RAG mode
                rag_prompt = f"""Based on the following context from the uploaded documents, please answer the user's question. If the context doesn't contain enough information to answer the question, please say so.

Context:
{context}

Question: {query}

Answer:"""
                system_message = self.rag_system_message
            
            # Format messages for OpenAI API
            messages = [
                {"role": "system", "content": system_message},
                {"role": "user", "content": rag_prompt}
            ]
            
            # Generate response
            response = await self.chat_model.arun(messages)
            return response
            
        except Exception as e:
            raise Exception(f"Error querying documents: {str(e)}")
    
    def query(self, query: str, k: int = 3, mode: str = "rag") -> str:
        """
        Synchronous wrapper for query_documents.
        
        Args:
            query: User query
            k: Number of chunks to retrieve
        
        Returns:
            Generated response based on retrieved chunks
        """
        try:
            # Search for relevant chunks
            relevant_chunks = self.search_relevant_chunks(query, k=k)
            
            if not relevant_chunks:
                return "I couldn't find any relevant information in the uploaded documents to answer your question."
            
            # Prepare context from chunks
            context_parts = []
            for chunk_data in relevant_chunks:
                chunk_text = chunk_data.get('chunk_text', '')
                if chunk_text:
                    context_parts.append(chunk_text)
            
            context = "\n\n".join(context_parts)
            
            # Create prompt based on mode
            if mode == "topic-explorer":
                # Topic Explorer mode with structured 3-part response
                rag_prompt = f"""You are a helpful and engaging study assistant for middle school students.
You will receive:
1. A student question
2. Context extracted from their study material (PDF, Word, or PowerPoint)

Your task:
- Use the context when possible, but also explain in a clear, simple way.
- Always return the answer in **three structured sections**:
  1. **Explanation** → Simple, student-friendly explanation of the topic.
  2. **Real-Life Example** → Show how the concept applies in daily life or something relatable.
  3. **Practice Activity** → Give a short challenge (question, drawing, or exercise) the student can do to practice.

Format your answer clearly with section headers.

---
Student Question:
{query}

Context:
{context}

Answer:"""
                system_message = "You are a helpful and engaging study assistant for middle school students."
            else:
                # Regular RAG mode
                rag_prompt = f"""Based on the following context from the uploaded documents, please answer the user's question. If the context doesn't contain enough information to answer the question, please say so.

Context:
{context}

Question: {query}

Answer:"""
                system_message = "You are a helpful assistant that answers questions based on provided context."
            
            # Format messages for OpenAI API
            messages = [
                {"role": "system", "content": system_message},
                {"role": "user", "content": rag_prompt}
            ]
            
            # Generate response synchronously
            response = self.chat_model.run(messages)
            return response
            
        except Exception as e:
            raise Exception(f"Error querying documents: {str(e)}")
    
    def get_document_info(self) -> Dict[str, Any]:
        """
        Get information about indexed documents.
        
        Returns:
            Dictionary containing document information
        """
        try:
            total_documents = len(self.documents)
            total_chunks = sum(doc_info.get('chunk_count', 0) for doc_info in self.documents.values())
            
            return {
                "total_documents": total_documents,
                "total_chunks": total_chunks,
                "documents": list(self.documents.keys()),
                "model_name": self.model_name
            }
        except Exception as e:
            return {
                "total_documents": 0,
                "total_chunks": 0,
                "documents": [],
                "model_name": self.model_name,
                "error": str(e)
            }


# Global RAG systems storage
rag_systems = {}

def get_or_create_rag_system(session_id: str, api_key: str, model_name: str = "gpt-4o-mini") -> RAGSystem:
    """Get existing or create new RAG system for session."""
    if session_id not in rag_systems:
        rag_systems[session_id] = RAGSystem(api_key=api_key, model_name=model_name)
    return rag_systems[session_id]